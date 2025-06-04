#!/usr/bin/env python3
# coding: utf-8

#Skriptet lager billetter med unik QR-kode basert på Excel-fil 

import os
import pandas as pd

import qrcode # type: ignore
import xlsxwriter
import openpyxl
from pptx import Presentation
from pptx.util import Inches
import sys
import win32com.client 
import PIL
import warnings
import time
import tkinter as tk
from tkinter import ttk
from tkinter import font
from tkinter import filedialog as fd
from tkinter.messagebox import showinfo
from tkinter import StringVar
from tkinter.filedialog import askopenfile
from tkinter.filedialog import askopenfilename
import matplotlib.pyplot as plt

#For å tilpasse programmet til nye arrangementer må bakgrunnsbildene for vanlig og ungdomsbillett oppdateres. Ta screenshot fra Powerpoint. 
#Videre må id, navn og billettype ligge i kolonnene [0, 10, 19].
#[19] må inneholder "Adult" for voksen/vanlig billett, "don't" for dem som ikke har bestilt. 

# Suppress specific warnings
warnings.filterwarnings("ignore", category=DeprecationWarning)
# Ignore all FutureWarnings
warnings.simplefilter(action='ignore', category=FutureWarning)

# Reading file within pyinstaller exe-file

if getattr(sys, 'frozen', False):
    application_path = os.path.dirname(sys.executable)
elif __file__:
    application_path = os.path.dirname(__file__)
    application_path = ''

def run_if_ready():
    if participants_first_name != "" and destination_folder != "" and qr_info !="" and ticket_template != "": 
        global button_create
        button_create=tk.Button(main_window, text="Lag billetter", fg="white", bg="red", command=create_tickets)
        button_create.place(x=434, y=550)
            

def my_path(path_name):
    #Return the appropriate path for data files based on execution context
    if getattr(sys, 'frozen', False):
        # running in a bundle
        return(os.path.join(sys._MEIPASS, path_name))

    else:
        # running live
        return path_name   

def on_closing(): #Tkinter
    main_window.destroy()    

global label_confirmed
global label_in
global label_ticket
global label_out
global label_progress
global button_create

def reset_process(window, new_text):
    #global qr_text
    
    for widget in window.winfo_children():
        if widget == qr_text:
            qr_text.delete('1.0', 'end')
        if widget == label_confirmed:
            label_confirmed.config(text=new_text)
    #global label_in
    label_in.config(text=new_text)
    #global label_ticket
    label_ticket.config(text=new_text)
    #global label_out
    label_out.config(text=new_text)
    label_progress.config(text=new_text)
    button_create.place_forget()

global qr_info
qr_info = ""

def get_qr_text():
    global qr_info    
    qr_info = qr_text.get("1.0", tk.END)  # Get text from line 1, 
    global label_confirmed
    label_confirmed = tk.Label(main_window, text="TEKST BEKREFTET", justify="left", font=('Open Sans', 8), bg="#6EC1E4", fg='white') 
    label_confirmed.place(x=275, y=482)
    print(qr_info)
    run_if_ready()

global pptx_mal_voksen
global pptx_mal_ungdom

global source_file
source_file = ""
global participants_form
participants_form = ""
global participants_first_name
participants_first_name = ""
global selected_ticket_file
selected_ticket_file = ""

def select_list_file(): #For tkinter
    global selected_participants
    global participants

    filetypes = (
    ('MS Excel-filer', '*.xlsx'),
    ('Csv-filer', '*.csv')
    )

    #global ticket_template
    participants = fd.askopenfilename(
        title='Velg fil',
        initialdir=os.path.expanduser("~"),
        filetypes=filetypes)
    
    if participants.endswith(".csv"):
        datasett=pd.read_csv(participants, sep=None, encoding='iso-8859-1', engine='python')
    if participants.endswith('.xlsx'):
        datasett=pd.read_excel(participants, engine='openpyxl')

    global participants_form
    participants_form = datasett.iloc[:, [0, 1, 2, 3]]
    participants_form.columns = ["Nr:", "Dato:", "Navn:", "Billett:"]

    selected_participants = 'Valgt deltakerliste: \n' + str(participants)
    global label_in
    label_in = tk.Label(main_window, text=selected_ticket_file,wraplength=500, justify="left", font=('Open Sans', 6), bg="#6EC1E4", fg='grey') 
    label_in.place(x=50, y=275)
    
    global participants_first_name
    participants_first_name=str(participants_form.loc[0,"Navn:"])
    run_if_ready()
        
global ticket_template
ticket_template = ""

def select_ticket_file(): #For tkinter
    filetypes = (
        ('Png-filer', '*.png'),
        ('JPG-filer', '*.jpg')
    )

    #global ticket_template
    global ticket_template
    ticket_template = fd.askopenfilename(
        title='Velg fil',
        initialdir=os.path.expanduser("~"),
        filetypes=filetypes)
    
    global selected_ticket_file
    selected_ticket_file = 'Valgt billettmal: \n' + str(ticket_template)
    global label_ticket
    label_ticket = tk.Label(main_window, text=selected_ticket_file,wraplength=500, justify="left", font=('Open Sans', 6), bg="#6EC1E4", fg='grey') 
    label_ticket.place(x=50, y=170)
    run_if_ready()


global destination_folder
destination_folder = ""

def select_folder(): #For tkinter
    global destination_folder
    destination_folder = fd.askdirectory(
        title="Velg mappe",
        initialdir=os.path.expanduser("~"))
    # Add a label to display the selected target folder
    
    chosen_folder = 'Lagringsmappe: \n' + str(destination_folder)
    global label_out
    label_out = tk.Label(main_window, text=chosen_folder,wraplength=500, justify="left", font=('Open Sans', 6), bg="#6EC1E4", fg='grey') 
    label_out.place(x=50, y=595)
    run_if_ready()

def create_tickets():
    global label_progress
    label_progress = tk.Label(main_window, text="Lager billetter ...", justify="left", font=('Open Sans', 8), bg="#6EC1E4", fg='white') 
    label_progress.place(x=275, y=550)

    global destination_folder
    folder=str(destination_folder)
    folder_qr = os.path.join(folder, 'QR-koder')
    global folder_bill
    folder_bill = os.path.join(folder, 'Billetter')

    if not os.path.exists(folder_qr):
        os.makedirs(folder_qr)
    if not os.path.exists(folder_bill):
        os.makedirs(folder_bill)

    global participants
    participants = participants_form.iloc[:, [0, 1, 2, 3]]
    participants.columns = ["Nr:", "Dato:", "Navn:", "Billett:"]

    participants= participants[participants['Billett:'].str.contains("don't")==False]#Tilpasning til SKRs påmeldingsliste
    participants['Registrert:'] = ""

    global qr_text
    qr_text_input=qr_info

    bestilling_filename = os.path.join(destination_folder, "Bestillingsliste.xlsx")
    participants = participants.sort_values('Nr:') 

    global pptx_cwd
    pptx_cwd = os.getcwd()
    global test_pptx
    test_pptx = 'test_pptx.pptx'
    
    global prs
    prs = Presentation()
    prs.slide_height = Inches(5.7)      
    prs.slide_width = Inches(10)
    blank_slide_layout = prs.slide_layouts[0]
    global slide
    slide = prs.slides.add_slide(blank_slide_layout)

    slide = prs.slides[0]

    for index, row in participants.iterrows():

        billettype = str(row.iloc[3]) #Sjekk ut variabelnavn
        
        navn = row[2]
        billettnummer = str(row[0])
        dato = str(row[1])

        # Convert string to pandas datetime object
        #date_obj = pd.to_datetime(dato)
        date_obj = pd.Timestamp(dato)

        # Extract day, month, and year
        day = date_obj.day

        month = date_obj.month
        year = date_obj.year
        datostr = str(str(day) + '.' + str(month) + '.' + str(year) + '_')

        participants.at[index, 'Dato:'] = datostr
        
        navn_under=navn.replace(" ", "_")
        if navn_under.endswith("_"):
            navn_under=navn_under[:-1]
        
        billettype = str(row[3])
        qr_filename= 'QR' + '_' + navn_under + '.png'

        pptx_filename = 'Billett_' + navn_under + '.pptx'
        pptx_filename = my_path(pptx_filename)
        global new_qr_file
        new_qr_file=os.path.join(folder_qr, qr_filename)
        #new_qr_file = Path(new_qr_file)
        global img_path
        img_path = new_qr_file
        global new_pptx_file
        new_pptx_file = os.path.join(folder_bill, pptx_filename)
        new_pptx_file = os.path.normpath(new_pptx_file)#Slettet Path fail
        #global filename_billett_pdf
        filename_billett_pdf = 'Billett_' + navn_under + '.pdf'
        global filenamefolder_billett_pdf
        filenamefolder_billett_pdf = os.path.join(folder_bill, filename_billett_pdf)
        filenamefolder_billett_pdf = os.path.normpath(filenamefolder_billett_pdf)#Slettet Path fail
        
        qr_text_ny = qr_text_input.replace('{navn}', navn)
        qr_text_ny = qr_text_ny.replace('{billettnummer}', billettnummer)
        qr_text_ny = qr_text_ny.replace('{billettype}', billettype)
        qr_text_ny = qr_text_ny.replace('{dato}', dato)
        qr_text_ny = qr_text_ny.replace('{ls} ', '{ls}') #Fjerner uønsket mellomrom på starten av linje
        qr_text_ny = qr_text_ny.replace('{ls}', '\n')

        # Create qr code instance
        qr = qrcode.QRCode(
            version = 1,
            error_correction = qrcode.constants.ERROR_CORRECT_L,
            box_size = 10,
            border = 4,
        )

        # Add data
        qr.add_data(qr_text_ny)
        qr.make(fit=True)

        # Create an image from the QR Code instance
        img = qr.make_image(fill_color="black", back_color="white")

        # Save the QR code 
        img.save(new_qr_file)
        img.close()

        # Print a success message
        
        left = Inches (3.8)
        left_bacground = Inches(0)
        top_background = Inches(0)
        top = Inches(2.5)
        height = Inches(1.5)
        width   = Inches(1.5)
        width_background = Inches (10)

        global ticket_template
        pptx_mal = ticket_template
    
        #Sqr_filename = new_qr_file
        slide.shapes.add_picture(pptx_mal, left=left_bacground, top=top_background, width=width_background)
        
        slide.shapes.add_picture(new_qr_file, left=left, top=top, width=width, height=height)
        
        new_pptx_file = my_path(new_pptx_file)

        new_pptx_file = os.path.join(folder_bill,new_pptx_file)
        new_pptx_file = os.path.normpath(new_pptx_file)

        prs.save(new_pptx_file)

        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        #powerpoint.Visible = False

        presentation = powerpoint.Presentations.Open(new_pptx_file, WithWindow=False) #

        presentation.Slides[0].Export(filenamefolder_billett_pdf, "PDF")
        #presentation.SaveAs(filenamefolder_billett_pdf, 32)  # 32 is the format type for PDF

        presentation.Close()
        powerpoint.Quit()
        presentation = None
        powerpoint = None
    
        os.remove(new_pptx_file)
        print('PDF-billett laget for ' + navn)
        

    if os.path.exists(bestilling_filename):
        total_bestilling = pd.read_excel(bestilling_filename, engine='openpyxl')
        frames=[total_bestilling, participants]
        total_bestilling_new = pd.concat(frames)
        total_bestilling_new = total_bestilling_new.drop_duplicates(subset='Nr:')
    else:
        total_bestilling_new = participants

    #total_bestilling_new = total_bestilling_new.sort_values(by=['Nr'], ascending=False)
    total_bestilling_new = total_bestilling_new.sort_values(by=['Nr:'])
    engine = 'xlsxwriter'

    with pd.ExcelWriter(bestilling_filename, engine=engine) as writer:
        total_bestilling_new.to_excel(writer, sheet_name="Bestillingsliste", index = None, header=True)
                    
    reset_process(main_window, "")
    
#MAIN WINDOW
main_window = tk.Tk()
main_window.protocol("WM_DELETE_WINDOW", on_closing) #Lukker Python ved å lukke vinduet
main_window.tk.call('tk', 'scaling', 2.0)
main_window.config(bg="#6EC1E4")  # Frame with a custom hex

tkinter_header_font = font.Font(size=16)
tkinter_general_font = font.Font(size=8)

main_window.title("CORals QR Ticket generator")
main_window.geometry ("600x650")

label = tk.Label(main_window, text="Billettgenerator", font=tkinter_header_font, bg="#6EC1E4")
label.place(x=50, y=50)
label = tk.Label(main_window, text='NB! Krav til kolonnenr. og -navn:\nID=1, Dato=2, Navn=3, Billettype=4', font=tkinter_general_font, justify='left', fg= 'white', bg="#6EC1E4")
label.place(x=225, y=225)
label = tk.Label(main_window, text='NB: Formatet må være 4:3 (std. Powerpoint)', font=tkinter_general_font, justify='left', fg= 'white', bg="#6EC1E4")
label.place(x=225, y=130)

# Add a button to trigger file selection
button = tk.Button(main_window, text="Velg bilettmal" , command=select_ticket_file)
button.place(x=50, y=125)

button = tk.Button(main_window, text="Velg deltakerliste" , command=select_list_file)
button.place(x=50, y=225)

label = tk.Label(main_window, text='Oppgi informasjonstekst til QR-kode:' , font=tkinter_general_font, bg="#6EC1E4", fg='black')
label.place(x=50, y=330)
label = tk.Label(main_window, text='Eksempel: "Dato = {dato}, Billettnr = {billettnummer}, Navn = {navn}, Billettype = {billettype}, Linjeskift = {ls}."', justify='left', wraplength=500, font=tkinter_general_font, bg="#6EC1E4", fg='white')
label.place(x=50, y=355)

qr_text = tk.Text(main_window, height=3, width=41, wrap="word")  # wrap="word" ensures text wraps at word boundaries
qr_text.place(x=50, y=400)

button = tk.Button(main_window, text="Bekreft tekst", command=get_qr_text)
button.place(x=434, y=475)

button = tk.Button(main_window, text='Velg lagringsmappe', command=select_folder)
button.place(x=50, y=550)



main_window.mainloop()




